# -*- coding:utf-8 -*-
import json
import re
import asyncio
import requests
import io
import time
import threading
import atexit
import subprocess
import tempfile
import gc
import os
import fitz
from concurrent.futures import ThreadPoolExecutor
from flask import Flask, request, jsonify
import config
# ==========新增向量库导入【RAG知识库】==========
from vector_kb import rag_answer, add_archive_to_kb, remove_archive_from_kb, rebuild_kb

# 全局环境变量，限制底层库线程，防止CPU/内存爆炸
os.environ["OMP_NUM_THREADS"] = "2"
os.environ["MKL_NUM_THREADS"] = "2"
os.environ["FLAGS_allocator_strategy"] = "naive_best_fit"

app = Flask(__name__)
app.json.ensure_ascii = False
app.config['JSONIFY_MIMETYPE'] = "application/json; charset=utf-8"
app.debug = False
app.jinja_env.auto_reload = False

tenant_token = ""
token_expire_time = 0
TOKEN_LOCK = threading.Lock()

# ==========试题持久缓存【/list_test /answer id】==========
TEST_CACHE_FILE = "test_record_cache.json"
TEST_CACHE_TTL = config.CACHE_EXPIRE_SECONDS
user_last_paper = {}
cache_lock = threading.Lock()

def load_test_cache():
    """从磁盘加载试题缓存（程序启动时调用一次）"""
    if not os.path.exists(TEST_CACHE_FILE):
        return {}
    try:
        with open(TEST_CACHE_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception as e:
        print(f"⚠️试题缓存读取失败：{str(e)}")
        return {}

def save_test_cache(data):
    try:
        with open(TEST_CACHE_FILE, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
    except Exception as e:
        print(f"⚠️试题缓存保存失败：{str(e)}")

# 启动时一次性载入内存，后续读写都基于内存，避免每条指令都读磁盘
test_cache_data = load_test_cache()

def clean_expired_test_records():
    global user_last_paper
    now = time.time()
    modified = False
    with cache_lock:
        for open_id in list(test_cache_data.keys()):
            record_list = test_cache_data[open_id]
            new_records = []
            for rec in record_list:
                if now - rec["create_time"] < TEST_CACHE_TTL:
                    new_records.append(rec)
                else:
                    modified = True
            test_cache_data[open_id] = new_records
            if len(test_cache_data[open_id]) == 0:
                del test_cache_data[open_id]
        if modified:
            save_test_cache(test_cache_data)
        user_last_paper = test_cache_data

def add_test_record(open_id: str, archive_id: int, question: str, answer: str):
    now = time.time()
    new_record = {
        "archive_id": archive_id,
        "question": question,
        "answer": answer,
        "create_time": now
    }
    with cache_lock:
        if open_id not in test_cache_data:
            test_cache_data[open_id] = []
        test_cache_data[open_id].append(new_record)
        save_test_cache(test_cache_data)
        global user_last_paper
        user_last_paper = test_cache_data

def get_latest_by_archive_id(open_id: str, target_archive_id: int):
    with cache_lock:
        records = test_cache_data.get(open_id, [])
        for rec in reversed(records):
            if rec["archive_id"] == target_archive_id:
                return rec
        return None

def get_latest_all_record(open_id: str):
    with cache_lock:
        records = test_cache_data.get(open_id, [])
        if not records:
            return None
        return records[-1]

def list_all_user_test_records(open_id: str):
    clean_expired_test_records()
    with cache_lock:
        return test_cache_data.get(open_id, [])

processed_event = {}
event_lock = threading.Lock()

last_pdf_bytes = None
last_pdf_name = ""
pdf_cache_lock = threading.Lock()

# 全局线程池，严格限制并发
executor = ThreadPoolExecutor(max_workers=config.MAX_WORKERS, thread_name_prefix="feishu_task")

def shutdown_executor():
    print("🛑正在关闭任务线程池...")
    executor.shutdown(wait=True, cancel_futures=True)
atexit.register(shutdown_executor)

def split_long_text(text: str, max_chars=4200) -> list[str]:
    chunks = []
    start = 0
    text_len = len(text)
    while start < text_len:
        end = start + max_chars
        cut_pos = text.rfind("\n", start, end)
        if cut_pos == -1 or cut_pos <= start:
            cut_pos = end
        chunks.append(text[start:cut_pos])
        start = cut_pos
    return chunks

def send_long_msg(open_id, full_text):
    parts = split_long_text(full_text)
    for idx, seg in enumerate(parts):
        send_msg(open_id, seg)
        time.sleep(0.25)

def clean_expired_cache():
    now = time.time()
    expired_event = []
    with event_lock:
        for eid, ts in processed_event.items():
            if now - ts > 300:
                expired_event.append(eid)
        for eid in expired_event:
            del processed_event[eid]
    if expired_event:
        print(f"🧹清理过期事件id数量：{len(expired_event)}")

def cache_clean_loop():
    while True:
        time.sleep(config.CLEAN_INTERVAL)
        try:
            clean_expired_cache()
            clean_expired_test_records()
            # 周期性垃圾回收
            gc.collect()
        except Exception as e:
            print("缓存清理异常：", str(e))

threading.Thread(target=cache_clean_loop, daemon=True).start()
print("⏱️试题缓存自动清理线程已启动，试题30分钟自动过期")

def get_valid_token():
    global tenant_token, token_expire_time
    now = time.time()
    if now >= token_expire_time - 30:
        with TOKEN_LOCK:
            if now >= token_expire_time - 30:
                url = "https://open.feishu.cn/open-apis/auth/v3/tenant_access_token/internal"
                payload = {
                    "app_id": config.FEISHU_APP_ID,
                    "app_secret": config.FEISHU_APP_SECRET
                }
                try:
                    with requests.post(url, json=payload, timeout=10) as resp:
                        res_data = resp.json()
                        if res_data.get("code") != 0:
                            raise Exception(f"获取token失败：{res_data.get('msg')}")
                        tenant_token = res_data["tenant_access_token"]
                        token_expire_time = now + res_data["expire"]
                        print("✅ 成功刷新tenant_access_token")
                except Exception as e:
                    print("❌ 获取Token异常：", str(e))
                    return None
    return tenant_token

def send_msg(open_id, text):
    token = get_valid_token()
    if not token:
        print("无法发送消息，Token为空")
        return None
    url = "https://open.feishu.cn/open-apis/im/v1/messages"
    headers = {"Authorization": f"Bearer {token}", "Content-Type":"application/json"}
    payload = {
        "receive_id": open_id,
        "msg_type": "text",
        "content": json.dumps({"text": text})
    }
    params = {"receive_id_type": "open_id"}
    try:
        with requests.post(url, headers=headers, params=params, json=payload, timeout=10) as resp:
            result = resp.json()
            if result.get("code") != 0:
                print(f"⚠️消息发送失败 返回：{result}")
            return result
    except Exception as e:
        print(f"❌调用发送消息接口异常：{str(e)}")
        return None

def send_interactive_card(open_id, card: dict):
    """发送飞书互动卡片（带按钮）"""
    token = get_valid_token()
    if not token:
        print("无法发送互动卡片，Token为空")
        return None
    url = "https://open.feishu.cn/open-apis/im/v1/messages"
    headers = {"Authorization": f"Bearer {token}", "Content-Type": "application/json"}
    payload = {
        "receive_id": open_id,
        "msg_type": "interactive",
        "content": json.dumps(card)
    }
    params = {"receive_id_type": "open_id"}
    try:
        with requests.post(url, headers=headers, params=params, json=payload, timeout=10) as resp:
            result = resp.json()
            if result.get("code") != 0:
                print(f"⚠️互动卡片发送失败 返回：{result}")
            return result
    except Exception as e:
        print(f"❌调用发送互动卡片接口异常：{str(e)}")
        return None

def build_help_text() -> str:
    """完整指令清单文本"""
    return """🤖可用指令清单：
/test 科目 内容 ｜ /test id 归档ID 生成自测题
/answer 查看最新答案 ｜ /answer id 归档ID 指定习题答案
/list_test 列出最近试题记录
/plan id 归档ID [days N] 🆕根据归档文档生成完整中长期学习计划
/daily id 归档ID [days N] 🆕生成每日学习任务
/cards id 归档ID 🆕把归档文档提炼成知识点和背诵卡片
/done id 归档ID day N 🆕打卡完成第N天学习任务
/progress id 归档ID 🆕查看学习进度
/save 科目 知识点 手动归档最近文件
/del id 数字 删除归档文档（同步清理向量库）
/rebuild_kb 全量重建向量知识库
/polish 德语/英语 文本 ｜ /polish id 归档ID ✍️润色修改德语/英语文本
/tip 打开互动菜单
上传PDF/DOC/DOCX/PPTX → 自动归档+自动出题+自动入库向量库
💡直接发送文字问题（不带/），自动检索本地文档进行问答
"""

def build_archive_list_text() -> str:
    """归档清单文本"""
    from archive_db import get_all_archive_summary
    archive_info = get_all_archive_summary()
    if isinstance(archive_info, str):
        return archive_info
    archive_lines = ["📂归档清单："]
    if isinstance(archive_info, list):
        for item in archive_info:
            if isinstance(item, dict):
                aid = item.get("id", "")
                subj = item.get("subject", "")
                fname = item.get("filename", "")
                archive_lines.append(f"ID:{aid} | {subj} | {fname}")
            else:
                print(f"[警告]归档列表非法数据，跳过，类型:{type(item)},内容:{item}")
    else:
        archive_lines.append("⚠️归档数据读取异常")
    return "\n".join(archive_lines)

def build_menu_card() -> dict:
    """互动菜单卡片（/tip 回复）"""
    return {
        "config": {"wide_screen_mode": True},
        "header": {"title": {"tag": "plain_text", "content": "📚 学习助手菜单"}, "template": "blue"},
        "elements": [
            {"tag": "div", "text": {"tag": "lark_md", "content": "点击按钮即可执行，选择你需要的功能："}},
            {"tag": "action", "actions": [
                {"tag": "button", "text": {"tag": "plain_text", "content": "📋 指令清单"}, "value": {"cmd": "help"}},
                {"tag": "button", "text": {"tag": "plain_text", "content": "📝 试题记录"}, "value": {"cmd": "list_test"}}
            ]},
            {"tag": "action", "actions": [
                {"tag": "button", "text": {"tag": "plain_text", "content": "🗂 归档清单"}, "value": {"cmd": "archives"}},
                {"tag": "button", "text": {"tag": "plain_text", "content": "✍️ 文本润色"}, "value": {"cmd": "polish"}}
            ]},
            {"tag": "action", "actions": [
                {"tag": "button", "text": {"tag": "plain_text", "content": "📇 背诵卡片"}, "value": {"cmd": "cards"}},
                {"tag": "button", "text": {"tag": "plain_text", "content": "🗑 重建知识库"}, "value": {"cmd": "rebuild_kb"}}
            ]}
        ]
    }

def build_menu_button_card() -> dict:
    """按钮处理完成后出现的“菜单”入口小卡片（点按钮才展开完整菜单）"""
    return {
        "config": {"wide_screen_mode": True},
        "header": {"title": {"tag": "plain_text", "content": "✅ 已处理完成"}, "template": "green"},
        "elements": [
            {"tag": "div", "text": {"tag": "lark_md", "content": "需要继续操作的话，点击下方按钮打开完整菜单："}},
            {"tag": "action", "actions": [
                {"tag": "button", "text": {"tag": "plain_text", "content": "📚 学习助手菜单"}, "value": {"cmd": "menu"}}
            ]}
        ]
    }

def format_msg(raw_text: str) -> str:
    text = raw_text.replace("\\\\(", "$")
    text = text.replace("\\\\)", "$")
    text = text.replace("\\\\begin", "\\begin")
    text = text.replace("\\\\end", "\\end")
    text = text.replace("\\\\\\\\", "\\\\")
    return text

def auto_extract_archive_info(pdf_text: str) -> dict:
    from llm_summary import llm_request
    short_text = pdf_text[:3000]
    prompt = f"""
你是专业课程资料归档助手，严格遵守下面所有规则，**只输出纯JSON字符串**，禁止输出任何前置说明、注释、markdown、思考过程、多余换行。
{{"subject":"科目名称","keypoint":"文档核心知识点概括，20～45字"}}
文档片段：
{short_text}
"""
    resp = llm_request(prompt)
    json_match = re.search(r"\{.*\}", resp, re.DOTALL)
    if not json_match:
        return {"subject":"未知科目","keypoint":"未识别知识点"}
    try:
        data = json.loads(json_match.group())
        subject = str(data.get("subject", "未知科目")).strip()
        kp = str(data.get("keypoint", "未识别知识点")).strip()
        return {"subject": subject, "keypoint": kp}
    except Exception:
        return {"subject":"未知科目","keypoint":"未识别知识点"}

def ai_simplify_filename(raw_name: str, subject: str) -> str:
    from llm_summary import llm_request
    prompt = f"""精简文件名，只输出名称，不要任何解释。
原始名称：{raw_name}
科目：{subject}
去除.pdf、副本、扫描版、水印、多余括号，控制20字内。
"""
    try:
        short_name = llm_request(prompt).strip().replace("\n","")
        if len(short_name) > 30:
            raise Exception("过长")
        return short_name
    except Exception:
        fallback = re.sub(r"[（(].*?[）)]", "", raw_name)
        fallback = re.sub(r"\.(pdf|docx|doc|pptx)","", fallback)
        return fallback.strip()

def clean_document_text(raw_text: str) -> str:
    raw_text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', "", raw_text)
    # 英文单词跨行时补空格，避免粘连成错误单词
    raw_text = re.sub(r"([A-Za-z])\n(?=[A-Za-z])", r"\1 ", raw_text)
    raw_text = re.sub(r"([^。！？；：\n])\n", r"\1", raw_text)
    raw_text = re.sub(r"\n{2,}", "\n\n", raw_text)
    raw_text = re.sub(r"\s+", " ", raw_text)
    raw_text = re.sub(r"第\s*\d+\s*页\s*/?\d*", "", raw_text)
    return raw_text.strip()

_ocr_engine = None
def ocr_pdf_stream(pdf_bytes: bytes) -> str:
    global _ocr_engine
    try:
        from paddleocr import PaddleOCR
        if _ocr_engine is None:
            print("⚠️初始化OCR引擎")
            # 内存优化配置，关闭无用模块，降低常驻内存
            _ocr_engine = PaddleOCR(
                use_angle_cls=True,
                lang="ch",
                use_gpu=False,
                use_doc_orientation_classify=False,
                use_doc_unwarping=False,
                use_textline_orientation=False,
                enable_memory_optim=True,
                rec_batch_num=1
            )
    except Exception as e:
        print(f"OCR加载失败：{e}")
        return ""
    doc = None
    try:
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        ocr_result = []
        for page in doc:
            pix = page.get_pixmap(dpi=200) # 降低分辨率减少内存占用
            img_data = pix.tobytes("png")
            pix = None # 主动释放pixmap
            res = _ocr_engine.ocr(img_data, cls=True)
            lines = []
            if res and res[0]:
                for line_info in res[0]:
                    lines.append(line_info[1][0])
            ocr_result.append("\n".join(lines))
        full_text = clean_document_text("\n\n".join(ocr_result))
        return full_text
    finally:
        if doc is not None:
            doc.close()
        fitz.TOOLS.store_shrink(100)
        gc.collect()

def extract_pdf_text(message_id:str, file_key:str):
    token = get_valid_token()
    url = f"https://open.feishu.cn/open-apis/im/v1/messages/{message_id}/resources/{file_key}?type=file"
    headers = {"Authorization": f"Bearer {token}"}
    max_retry=2
    file_bytes = b""
    for retry_count in range(max_retry+1):
        try:
            with requests.get(url, headers=headers, allow_redirects=True, timeout=15) as resp:
                file_bytes = resp.content
            print(f"✅文件下载完成，字节大小：{len(file_bytes)}")
            if len(file_bytes)<200 and retry_count<max_retry:
                time.sleep(1)
                continue
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            full_text = ""
            for page in doc:
                blocks = page.get_text("blocks", sort=True)
                page_text = "\n".join([b[4].strip() for b in blocks if b[4].strip()])
                full_text += page_text + "\n\n"
            doc.close()
            fitz.TOOLS.store_shrink(100)
            full_text = clean_document_text(full_text)
            print(f"✅PDF原生文本提取完毕，文本长度：{len(full_text)}")
            if len(full_text.strip())<200:
                print("⚠️文本过少，启动OCR识别")
                full_text = ocr_pdf_stream(file_bytes)
            return full_text, file_bytes
        except Exception as e:
            print(f"⚠️PDF下载/解析重试 {retry_count}, err:{str(e)}")
            if retry_count<max_retry:
                time.sleep(1)
                continue
            raise Exception(str(e))

def extract_docx_text(message_id:str, file_key:str):
    token = get_valid_token()
    url = f"https://open.feishu.cn/open-apis/im/v1/messages/{message_id}/resources/{file_key}?type=file"
    headers = {"Authorization": f"Bearer {token}"}
    with requests.get(url, headers=headers, allow_redirects=True, timeout=15) as resp:
        file_bytes = resp.content
    from docx import Document
    with io.BytesIO(file_bytes) as stream:
        doc = Document(stream)
        text_list = []
        for p in doc.paragraphs:
            if p.text.strip():
                text_list.append(p.text.strip())
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text_list.append(cell.text.strip())
    full_text = clean_document_text("\n".join(text_list))
    return full_text, file_bytes

def extract_pptx_text(message_id:str, file_key:str):
    token = get_valid_token()
    url = f"https://open.feishu.cn/open-apis/im/v1/messages/{message_id}/resources/{file_key}?type=file"
    headers = {"Authorization": f"Bearer {token}"}
    with requests.get(url, headers=headers, allow_redirects=True, timeout=15) as resp:
        file_bytes = resp.content
    from pptx import Presentation
    with io.BytesIO(file_bytes) as stream:
        prs = Presentation(stream)
        text_list = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape,"text") and shape.text.strip():
                    text_list.append(shape.text.strip())
    full_text = clean_document_text("\n\n".join(text_list))
    return full_text, file_bytes

def extract_old_doc_text(file_bytes: bytes) -> str:
    try:
        with tempfile.NamedTemporaryFile(suffix=".doc", delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        result = subprocess.run(["antiword", tmp_path], capture_output=True, text=True, timeout=30)
        os.unlink(tmp_path)
        return clean_document_text(result.stdout)
    except Exception:
        return ""

# =====================【学习计划相关函数】=====================
def generate_study_plan(doc_content: str, subject: str):
    from llm_summary import llm_request
    prompt = f"""
你是学习规划助手。基于下面课程文档内容，生成一份结构化学习计划。
要求：
1. 划分学习阶段（预习→精读→习题→复盘）
2. 合理分配每日任务，建议3~7天学习周期
3. 标出重点、难点、自测方式
4. 排版清晰，不要多余废话
科目：{subject}
文档内容：
{doc_content[:6000]}
"""
    res = llm_request(prompt)
    return format_msg(res)

def split_plan_to_daily_tasks(full_plan: str, subject: str, days: int = 5):
    from llm_summary import llm_request
    prompt = f"""
你是学习拆解助手。
已有完整中长期学习规划，均衡拆分为【{days}天每日学习清单】
要求：
1. 每一天明确：学习内容、重点、自测任务、预估耗时
2. 难度循序渐进，前面预习精读，后面刷题复盘
3. 严格格式：
## Day1
任务内容：xxx
重点：xxx
自测：xxx

## Day2
任务内容：xxx
重点：xxx
自测：xxx
科目：{subject}
整体学习规划：
{full_plan[:6000]}
最后不要额外总结，只输出结构化内容。
"""
    resp = llm_request(prompt)
    return format_msg(resp)

from llm_summary import extract_task_list

# =====================【消息主逻辑】=====================
def process_message_task(event_data):
    global last_pdf_bytes, last_pdf_name
    receive_id = ""
    try:
        event = event_data.get("event", {})
        message = event.get("message", {})
        msg_type = message.get("message_type")
        sender_open_id = event.get("sender", {}).get("sender_id", {}).get("open_id")
        receive_id = sender_open_id
        print(f"【消息调试】msg_type={msg_type}, sender_open_id={sender_open_id}")

        if sender_open_id != config.ALLOW_OPEN_ID:
            send_msg(receive_id, "权限不足，仅本人可用")
            return

        if msg_type == "text":
            content = json.loads(message["content"])["text"].strip()
            print("收到指令：", content)

            if content.startswith("/test"):
                clean_expired_test_records()
                parts = content.split(maxsplit=2)
                archive_id = 0
                if len(parts)>=3 and parts[1].lower()=="id":
                    try:
                        archive_id=int(parts[2])
                    except ValueError:
                        send_msg(receive_id, "用法：/test id 归档ID")
                        return
                    from archive_db import get_archive_by_id
                    row = get_archive_by_id(archive_id)
                    if not row:
                        send_msg(receive_id, f"❌找不到归档ID {archive_id}")
                        return
                    subject = row["subject"]
                    text_content = row["file_text"]
                    send_msg(receive_id, f"📖加载归档ID:{archive_id}【{subject}】生成习题")
                else:
                    if len(parts)<3:
                        send_msg(receive_id, """📝出题方式：
1) /test 科目 文本
2) /test id 归档ID
/list_test 查看试题记录
""")
                        return
                    subject = parts[1]
                    text_content = parts[2]
                from llm_summary import generate_test_questions
                q,a = generate_test_questions(text_content, subject,10)
                cq = format_msg(q)
                ca = format_msg(a)
                add_test_record(sender_open_id, archive_id, cq, ca)
                send_long_msg(receive_id, f"{cq}\n💡发送 /answer 查看答案")

            elif content.startswith("/answer"):
                args = content.split(maxsplit=2)
                target_aid = None
                if len(args)>=3 and args[1].lower()=="id":
                    try:
                        target_aid=int(args[2])
                    except ValueError:
                        send_msg(receive_id, "格式：/answer id 归档ID")
                        return
                if target_aid is not None:
                    rec = get_latest_by_archive_id(sender_open_id, target_aid)
                else:
                    rec = get_latest_all_record(sender_open_id)
                if not rec:
                    send_msg(receive_id, "暂无试题记录，请先执行/test")
                    return
                send_long_msg(receive_id, f"【参考答案】\n{rec['answer']}")

            elif content == "/list_test":
                records = list_all_user_test_records(sender_open_id)
                if not records:
                    send_msg(receive_id, "📋暂无试题记录")
                    return
                lines = ["📋试题清单（30分钟有效期）："]
                for idx,r in enumerate(records,1):
                    aid = r["archive_id"]
                    desc = f"归档ID={aid}" if aid !=0 else "临时文本出题"
                    lines.append(f"{idx}. {desc}")
                lines.append("\n/answer id 数字 查询对应习题答案")
                send_msg(receive_id, "\n".join(lines))

            elif content.startswith("/plan id"):
                parts = content.split()
                aid = None
                target_days = 5
                try:
                    aid = int(parts[2])
                    if len(parts) >=5 and parts[3].lower() == "days":
                        target_days = int(parts[4])
                        target_days = max(2, min(14, target_days))
                except Exception:
                    send_msg(receive_id, "📖用法：\n/plan id 归档ID\n/plan id 归档ID days 7\n示例：/plan id 5 days 7")
                    return
                from archive_db import get_archive_by_id
                row = get_archive_by_id(aid)
                if not row:
                    send_msg(receive_id, f"❌未找到归档ID={aid}")
                    return
                send_msg(receive_id, f"🤖正在生成【{target_days}天】学习方案，请稍候...")
                full_plan_text = generate_study_plan(row["file_text"], row["subject"])
                daily_task_text = split_plan_to_daily_tasks(full_plan_text, row["subject"], target_days)
                output = f"""📅【{row['subject']}学习总方案】
归档ID：{aid}
文档名称：{row['filename']}
周期：{target_days}天

====整体规划====
{full_plan_text}

====📆每日细化任务清单====
{daily_task_text}
"""
                send_long_msg(receive_id, output)

            elif content.startswith("/daily id"):
                parts = content.split()
                aid = None
                target_days = 5
                try:
                    aid = int(parts[2])
                    if len(parts) >=5 and parts[3].lower() == "days":
                        target_days = int(parts[4])
                        target_days = max(2, min(14, target_days))
                except Exception:
                    send_msg(receive_id, """📖用法：
/daily id 归档ID
/daily id 归档ID days 6
""")
                    return
                from archive_db import get_archive_by_id
                row = get_archive_by_id(aid)
                if not row:
                    send_msg(receive_id, f"❌未找到归档ID={aid}")
                    return
                send_msg(receive_id, f"🤖正在拆分{target_days}天每日学习任务...")
                full_plan_text = generate_study_plan(row["file_text"], row["subject"])
                daily_task_text = split_plan_to_daily_tasks(full_plan_text, row["subject"], target_days)
                task_list = extract_task_list(daily_task_text)

                from review_scheduler import save_daily_tasks
                save_daily_tasks(aid, row["subject"], task_list)

                output = f"""📆【{row['subject']}每日学习任务】
归档ID：{aid} ｜ {target_days}天周期
{daily_task_text}
💡打卡用法：/done id {aid} day 1 标记第1天完成
📊进度查询：/progress id {aid}
                """
                send_long_msg(receive_id, output)

            elif content.startswith("/cards id"):
                parts = content.split()
                try:
                    aid = int(parts[2])
                except Exception:
                    send_msg(receive_id, "📝用法：/cards id 归档ID\n示例：/cards id 3")
                    return
                from archive_db import get_archive_by_id
                row = get_archive_by_id(aid)
                if not row:
                    send_msg(receive_id, f"❌未找到归档ID={aid}")
                    return
                if not row["file_text"] or len(row["file_text"].strip()) < 20:
                    send_msg(receive_id, "该归档文档没有可用的文本内容")
                    return
                send_msg(receive_id, f"🤖正在把归档ID:{aid}【{row['subject']}】提炼成知识点和背诵卡片，请稍候...")
                from llm_summary import generate_memory_cards
                try:
                    cards = generate_memory_cards(row["file_text"], row["subject"])
                except Exception as e:
                    send_msg(receive_id, f"❌生成失败：{str(e)}")
                    return
                send_long_msg(receive_id, f"📚【{row['subject']}】知识点与背诵卡片\n{format_msg(cards)}")

            elif content.startswith("/done id"):
                parts = content.split()
                try:
                    aid = int(parts[2])
                    day_num = int(parts[4])
                except Exception:
                    send_msg(receive_id, "📝用法：/done id 归档ID day 天数\n示例：/done id 6 day 2")
                    return
                from review_scheduler import mark_task_finished
                ok = mark_task_finished(aid, day_num)
                if ok:
                    send_msg(receive_id, f"✅已标记归档ID:{aid} 第{day_num}天任务完成！")
                else:
                    send_msg(receive_id, f"❌未找到对应任务，检查归档ID或天数是否正确")

            elif content.startswith("/progress id"):
                parts = content.split()
                try:
                    aid = int(parts[2])
                except Exception:
                    send_msg(receive_id, "📝用法：/progress id 归档ID")
                    return
                from review_scheduler import get_archive_progress
                rows = get_archive_progress(aid)
                if not rows:
                    send_msg(receive_id, "⚠️暂无任务记录，请先执行 /daily id xxx 生成每日任务")
                    return
                total = len(rows)
                finished = sum(1 for r in rows if r["finished"] == 1)
                rate = f"{finished/total*100:.1f}%" if total>0 else "0%"
                msg_lines = [f"📊学习进度 归档ID:{aid}\n总任务：{total}天 | 已完成：{finished} | 完成率：{rate}\n"]
                for r in rows:
                    status = "✅已完成" if r["finished"] else "⏳待完成"
                    complete_day = r["complete_date"] if r["complete_date"] else "未打卡"
                    msg_lines.append(f"Day{r['day_no']} {status} | {complete_day}")
                send_long_msg(receive_id, "\n".join(msg_lines))

            elif content.startswith("/save"):
                parts = content.split(maxsplit=2)
                if len(parts)<3:
                    send_msg(receive_id, "/save 科目 知识点")
                    return
                subject = parts[1]
                kname = parts[2]
                clean_name = ai_simplify_filename(kname, subject)
                with pdf_cache_lock:
                    if not last_pdf_bytes:
                        send_msg(receive_id, "⚠️先上传文件再执行/save")
                        return
                    from archive_db import archive_file
                    # =========修复参数顺序 + 接收元组返回值==========
                    save_path, new_id = archive_file(subject, clean_name, last_pdf_bytes, last_pdf_name, "")
                    send_msg(receive_id,f"✅归档成功！ID={new_id}")

            elif content.startswith("/del"):
                body = content.removeprefix("/del").strip()
                parts = body.split(maxsplit=2)
                if len(parts)>=2 and parts[0].lower()=="id":
                    try:
                        target_id=int(parts[1])
                        from archive_db import delete_archive_by_id
                        deleted_name = delete_archive_by_id(target_id)
                        if not deleted_name:
                            send_msg(receive_id,"找不到该归档ID")
                            return
                        remove_archive_from_kb(target_id)
                        send_msg(receive_id,f"✅删除归档ID:{target_id} {deleted_name}")
                    except ValueError:
                        send_msg(receive_id,"用法 /del id 数字")
                    return
                if "|" not in body:
                    send_msg(receive_id,"/del id 数字 推荐使用")
                    return
                subj,fname = [x.strip() for x in body.split("|",maxsplit=1)]
                from archive_db import delete_archive_file
                res = delete_archive_file(subj,"",fname)
                send_msg(receive_id, res)

            elif content.startswith("/polish") or content.startswith("/改写"):
                body = content.removeprefix("/polish").removeprefix("/改写").strip()
                if not body:
                    send_msg(receive_id, """📝用法：
/polish 德语 你的文本
/polish 英语 你的文本
/polish 你的文本（自动识别语言）
/polish id 归档ID（润色归档文档）
可以在文本前附上修改要求，例如：
/polish 德语 改得更口语化一点：原文内容
""")
                    return
                first, _, rest = body.partition(" ")
                if first == "id":
                    try:
                        aid = int(rest.strip())
                    except ValueError:
                        send_msg(receive_id, "用法：/polish id 归档ID")
                        return
                    from archive_db import get_archive_by_id
                    row = get_archive_by_id(aid)
                    if not row:
                        send_msg(receive_id, f"❌找不到归档ID {aid}")
                        return
                    lang = ""
                    text = row["file_text"]
                    if not text or len(text.strip()) < 20:
                        send_msg(receive_id, "该归档文档没有可用的文本内容")
                        return
                elif first in ("德语", "德", "de", "英语", "英", "en"):
                    lang = first
                    text = rest.strip()
                else:
                    lang = ""
                    text = body
                if not text:
                    send_msg(receive_id, "请把需要修改的文本一起发给我")
                    return
                send_msg(receive_id, "🔄正在润色，请稍候...")
                from llm_summary import polish_text
                try:
                    result = polish_text(text, lang)
                except Exception as e:
                    send_msg(receive_id, f"❌润色失败：{str(e)}")
                    return
                send_long_msg(receive_id, f"✍️润色结果：\n{result}")

            elif content == "/rebuild_kb":
                send_msg(receive_id, "🔄开始重建全部向量知识库，耗时较长，请耐心等待...")
                try:
                    rebuild_kb()
                    send_msg(receive_id, "✅知识库重建完成！所有归档文档已载入向量库")
                except Exception as e:
                    import traceback
                    traceback.print_exc()
                    send_msg(receive_id, f"❌重建知识库失败：{str(e)}")

            elif content == "/tip":
                send_interactive_card(receive_id, build_menu_card())

            else:
                send_msg(receive_id, "🤖正在检索本地归档资料，请稍候...")
                try:
                    reply = rag_answer(content)
                    send_long_msg(receive_id, reply)
                except Exception as e:
                    import traceback
                    traceback.print_exc()
                    send_msg(receive_id, f"问答服务出错：{str(e)}")

        elif msg_type == "file":
            print("【消息调试】检测到文件消息，准备解析")
            try:
                content_data = json.loads(message["content"])
                file_key = content_data["file_key"]
                file_name = content_data["file_name"]
                message_id = message.get("message_id")
                print(f"【文件消息调试】文件名={file_name}, key={file_key}")
                suffix = file_name.lower()
                doc_text, file_bytes = "", b""
                if suffix.endswith(".pdf"):
                    doc_text, file_bytes = extract_pdf_text(message_id, file_key)
                elif suffix.endswith(".docx"):
                    doc_text, file_bytes = extract_docx_text(message_id, file_key)
                elif suffix.endswith(".pptx"):
                    doc_text, file_bytes = extract_pptx_text(message_id, file_key)
                elif suffix.endswith(".doc"):
                    send_msg(receive_id, "解析老式DOC中...")
                    token = get_valid_token()
                    url = f"https://open.feishu.cn/open-apis/im/v1/messages/{message_id}/resources/{file_key}?type=file"
                    headers = {"Authorization": f"Bearer {token}"}
                    with requests.get(url, headers=headers, timeout=15) as resp:
                        file_bytes = resp.content
                    doc_text = extract_old_doc_text(file_bytes)
                else:
                    send_msg(receive_id, "不支持该文件格式")
                    return

                if len(doc_text.strip())<20:
                    send_msg(receive_id, "文档文字过少，无法处理")
                    return
                with pdf_cache_lock:
                    last_pdf_bytes = file_bytes
                    last_pdf_name = file_name
                doc_text = doc_text[:config.MAX_LLM_CONTEXT]
                auto_info = auto_extract_archive_info(doc_text)
                subj = auto_info["subject"]
                short_name = ai_simplify_filename(file_name, subj)
                from archive_db import archive_file
                # =========核心修复：参数顺序修正 + 接收元组==========
                save_path, new_aid = archive_file(subj, short_name, file_bytes, file_name, doc_text)
                add_archive_to_kb(new_aid)
                notice = f"""📂AI自动归档完成
科目：{subj}
归档ID：{new_aid}
💡出题：/test id {new_aid}
💡生成学习计划：/plan id {new_aid}
💡直接提问文档内容（无需指令）进行知识库问答
"""
                send_msg(receive_id, notice)
                from llm_summary import generate_test_questions
                q,a = generate_test_questions(doc_text, subj,10)
                cq,ca = format_msg(q), format_msg(a)
                add_test_record(sender_open_id, new_aid, cq, ca)
                send_long_msg(receive_id, f"{cq}\n💡发送 /answer 查看答案")
                # 主动释放大内存变量
                doc_text = None
                file_bytes = None
                gc.collect()
            except Exception as e:
                import traceback
                err_info = traceback.format_exc()
                print("====文件处理异常堆栈====")
                print(err_info)
                send_msg(receive_id, f"文件解析失败：{str(e)}")
                return
    except Exception as e:
        import traceback
        print("====任务顶层异常堆栈====")
        print(traceback.format_exc())
        if receive_id:
            send_msg(receive_id, f"任务异常：{str(e)}")

def task_future_callback(future):
    try:
        future.result()
    except Exception:
        import traceback
        print("====线程池任务捕获异常====")
        print(traceback.format_exc())

def handle_card_action(open_id: str, cmd: str):
    """处理互动卡片按钮点击，处理完成后再次发送菜单方便继续操作"""
    try:
        if cmd == "help":
            send_long_msg(open_id, build_help_text())
        elif cmd == "archives":
            send_long_msg(open_id, build_archive_list_text())
        elif cmd == "list_test":
            records = list_all_user_test_records(open_id)
            if records:
                lines = ["📋试题清单（30分钟有效期）："]
                for idx, r in enumerate(records, 1):
                    aid = r["archive_id"]
                    desc = f"归档ID={aid}" if aid != 0 else "临时文本出题"
                    lines.append(f"{idx}. {desc}")
                lines.append("\n/answer id 数字 查询对应习题答案")
                send_msg(open_id, "\n".join(lines))
            else:
                send_msg(open_id, "📋暂无试题记录")
        elif cmd == "polish":
            send_msg(open_id, """📝润色用法：
/polish 德语 你的文本
/polish 英语 你的文本
/polish id 归档ID
/改写 你的文本
可在文本前附修改要求，例如：/polish 德语 更口语化：xxx""")
        elif cmd == "cards":
            send_msg(open_id, """📇背诵卡片用法：
/cards id 归档ID
示例：/cards id 3
我会把该归档文档提炼成核心知识点和背诵卡片（中文讲解，德语/英语内容配原文+释义+例句）""")
        elif cmd == "rebuild_kb":
            send_msg(open_id, "🔄开始重建全部向量知识库，耗时较长，请耐心等待...")
            try:
                rebuild_kb()
                send_msg(open_id, "✅知识库重建完成！所有归档文档已载入向量库")
            except Exception as e:
                send_msg(open_id, f"❌重建知识库失败：{str(e)}")
        elif cmd == "menu":
            send_interactive_card(open_id, build_menu_card())
            return
        elif cmd == "tip":
            send_interactive_card(open_id, build_menu_card())
            return
        else:
            send_msg(open_id, f"未知操作：{cmd}")
            return
        # 处理完成后发送“菜单入口”小卡片，点按钮再展开完整菜单
        send_interactive_card(open_id, build_menu_button_card())
    except Exception as e:
        import traceback
        traceback.print_exc()
        send_msg(open_id, f"❌操作失败：{str(e)}")

def handle_card_action_event(data: dict):
    """解析卡片点击回调并异步执行"""
    event = data.get("event", {})
    action = event.get("action", {})
    value = action.get("value", {}) or {}
    cmd = value.get("cmd", "")
    operator = event.get("operator", {})
    open_id = (
        operator.get("operator_id", {}).get("open_id")
        or operator.get("open_id")
        or action.get("open_id")
        or ""
    )
    if not open_id:
        return jsonify({"toast": {"type": "error", "content": "无法识别操作者"}})
    if not cmd:
        return jsonify({"toast": {"type": "error", "content": "无效操作"}})
    if open_id != config.ALLOW_OPEN_ID:
        return jsonify({"toast": {"type": "error", "content": "权限不足"}})
    executor.submit(handle_card_action, open_id, cmd)
    return jsonify({"toast": {"type": "info", "content": "正在处理，请稍候..."}})

@app.route("/feishu/card_callback", methods=["POST"])
def card_callback():
    try:
        return handle_card_action_event(request.get_json(silent=True) or {})
    except Exception as e:
        print(f"卡片回调异常：{e}")
        return jsonify({})

@app.route("/feishu/callback", methods=["POST"])
def callback():
    try:
        data = request.get_json()
        if "challenge" in data:
            return jsonify({"challenge": data["challenge"]})
        # 卡片点击事件（如果走事件订阅方式推送）
        if data.get("header", {}).get("event_type") == "card.action.trigger":
            return handle_card_action_event(data)
        event_id = data.get("header", {}).get("event_id")
        with event_lock:
            if event_id in processed_event:
                return jsonify({"code":0})
            processed_event[event_id] = time.time()
        fut = executor.submit(process_message_task, data)
        fut.add_done_callback(task_future_callback)
        return jsonify({"code":0})
    except Exception as e:
        print(e)
        return jsonify({"code":-1}),400

def start_feishu_bot():
    print("🤖飞书机器人启动【RAG知识库问答 /plan /daily /done /progress 全套功能｜内存优化版】")
    # 注意：开发服务器长期运行容易内存泄漏，生产环境建议改用gunicorn
    app.run(host="0.0.0.0", port=config.FLASK_PORT, debug=False, threaded=True, use_reloader=False)

if __name__ == "__main__":
    from archive_db import init_db
    init_db()
    start_feishu_bot()
