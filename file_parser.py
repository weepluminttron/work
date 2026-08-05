# -*- coding: utf-8 -*-
"""文件解析公共模块：网页版和飞书版共用"""
import io
import os
import re
import gc
import base64
import tempfile
import subprocess


def clean_document_text(raw_text: str) -> str:
    raw_text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', "", raw_text)
    raw_text = re.sub(r"([A-Za-z])\n(?=[A-Za-z])", r"\1 ", raw_text)
    raw_text = re.sub(r"([^。！？；：\n])\n", r"\1", raw_text)
    raw_text = re.sub(r"\n{2,}", "\n\n", raw_text)
    raw_text = re.sub(r"\s+", " ", raw_text)
    raw_text = re.sub(r"第\s*\d+\s*页\s*/?\d*", "", raw_text)
    return raw_text.strip()


_ocr_engine = None
_ocr_init_attempted = False


def _init_ocr_engine():
    """初始化本地 PaddleOCR 引擎，自动兼容不同版本（2.x/3.x 参数不同）"""
    global _ocr_engine, _ocr_init_attempted
    if _ocr_engine is not None:
        return _ocr_engine
    if _ocr_init_attempted:
        return None
    _ocr_init_attempted = True
    from paddleocr import PaddleOCR
    param_sets = [
        {"lang": "ch", "use_gpu": False, "enable_mkldnn": False},
        {"lang": "ch", "use_gpu": False},
        {"lang": "ch", "use_gpu": False, "use_angle_cls": True},
        {
            "lang": "ch", "use_gpu": False, "use_angle_cls": True,
            "use_doc_orientation_classify": False, "use_doc_unwarping": False,
            "enable_memory_optim": True, "rec_batch_num": 1,
        },
        {"lang": "ch"},
    ]
    last_err = None
    for kw in param_sets:
        try:
            _ocr_engine = PaddleOCR(**kw)
            print("✅本地OCR引擎初始化成功")
            return _ocr_engine
        except Exception as e:
            last_err = e
            continue
    print(f"⚠️本地OCR引擎初始化失败：{last_err}")
    return None


def _ocr_call(img) -> list:
    """兼容 PaddleOCR 2.x ocr() 与 3.x predict()，返回识别出的每行文字"""
    if _ocr_engine is None:
        return []
    res = None
    try:
        res = _ocr_engine.ocr(img)
    except Exception as e1:
        try:
            res = _ocr_engine.predict(img)
        except Exception as e2:
            print(f"⚠️本地OCR调用失败：ocr={e1}，predict={e2}")
            print("💡如果是 PaddlePaddle 3.3.0 的 oneDNN 兼容问题，请在服务器执行：pip install paddlepaddle==3.2.2")
            return []
    lines = []
    if not res:
        return lines
    for item in res:
        if isinstance(item, dict):
            rec_texts = item.get("rec_texts") or []
            for t in rec_texts:
                if t:
                    lines.append(str(t))
        elif hasattr(item, "rec_texts"):
            for t in item.rec_texts or []:
                if t:
                    lines.append(str(t))
        else:
            for line_info in item or []:
                if line_info and len(line_info) >= 2 and line_info[1] and line_info[1][0]:
                    lines.append(str(line_info[1][0]))
    return lines


def ocr_pdf_stream(pdf_bytes: bytes) -> str:
    """扫描版 PDF OCR（可选，未安装 paddleocr 时返回空）"""
    import fitz
    try:
        if _init_ocr_engine() is None:
            return ""
    except Exception as e:
        print(f"OCR加载失败：{e}")
        return ""
    doc = None
    try:
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        ocr_result = []
        for page in doc:
            pix = page.get_pixmap(dpi=200)
            img_data = pix.tobytes("png")
            pix = None
            lines = _ocr_call(img_data)
            ocr_result.append("\n".join(lines))
        return clean_document_text("\n\n".join(ocr_result))
    finally:
        if doc is not None:
            doc.close()
        fitz.TOOLS.store_shrink(100)
        gc.collect()


def extract_pdf_bytes_text(file_bytes: bytes) -> str:
    import fitz
    doc = None
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        full_text = ""
        for page in doc:
            blocks = page.get_text("blocks", sort=True)
            page_text = "\n".join([b[4].strip() for b in blocks if b[4].strip()])
            full_text += page_text + "\n\n"
        doc.close()
        doc = None
        fitz.TOOLS.store_shrink(100)
        full_text = clean_document_text(full_text)
        if len(full_text.strip()) < 200:
            print("⚠️文本过少，启动OCR识别")
            full_text = ocr_pdf_stream(file_bytes)
        return full_text
    finally:
        if doc is not None:
            doc.close()


def extract_docx_bytes_text(file_bytes: bytes) -> str:
    from docx import Document
    with io.BytesIO(file_bytes) as stream:
        doc = Document(stream)
        text_list = []
        for p in doc.paragraphs:
            if p.text.strip():
                text_list.append(p.text.strip())
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text_list.append(cell.text.strip())
    return clean_document_text("\n".join(text_list))


def extract_pptx_bytes_text(file_bytes: bytes) -> str:
    from pptx import Presentation
    with io.BytesIO(file_bytes) as stream:
        prs = Presentation(stream)
        text_list = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_list.append(shape.text.strip())
    return clean_document_text("\n\n".join(text_list))


def extract_old_doc_bytes_text(file_bytes: bytes) -> str:
    try:
        with tempfile.NamedTemporaryFile(suffix=".doc", delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        result = subprocess.run(["antiword", tmp_path], capture_output=True, text=True, timeout=30)
        os.unlink(tmp_path)
        return clean_document_text(result.stdout)
    except Exception:
        return ""


last_image_error = ""


def _friendly_api_error(status_code: int, resp_text: str) -> str:
    """把视觉接口的错误转换成用户能看懂的中文提示"""
    if status_code == 401:
        return "硅基流动 API 密钥无效（401），请检查服务器 .env 中的 EMB_API_KEY"
    if status_code == 402:
        return "硅基流动账户余额不足（402），请到控制台充值后重试"
    if status_code == 400:
        body = resp_text[:200] if resp_text else ""
        return f"视觉模型不可用（400）：{body}"
    return f"视觉接口返回异常（{status_code}）：{resp_text[:200]}"


last_local_ocr_error = ""


def _ocr_image_bytes_local(file_bytes: bytes, mime: str = "image/png") -> str:
    """使用本地 PaddleOCR 识别图片文字（免费，不需要 API 余额）"""
    global last_local_ocr_error
    last_local_ocr_error = ""
    try:
        import numpy as np
        from PIL import Image
        if _init_ocr_engine() is None:
            last_local_ocr_error = "本地OCR引擎初始化失败，请查看日志"
            return ""
        img = Image.open(io.BytesIO(file_bytes))
        arr = np.array(img.convert("RGB"))
        print("⏳本地OCR识别中（首次识别可能较慢）...")
        lines = _ocr_call(arr)
        if not lines:
            last_local_ocr_error = "本地OCR未识别出文字，图片可能不够清晰或为纯照片"
            return ""
        text = clean_document_text("\n".join(lines))
        print(f"✅本地OCR识别完成：{len(lines)} 行文字")
        return text
    except ImportError:
        last_local_ocr_error = "服务器未安装本地OCR（可执行 pip install paddleocr paddlepaddle）"
        print("⚠️本地OCR未安装：如需免费离线识别，请执行 pip install paddleocr paddlepaddle")
        return ""
    except Exception as e:
        last_local_ocr_error = f"本地OCR识别失败：{e}"
        print(f"⚠️本地OCR识别失败：{e}")
        return ""


def extract_image_text(file_bytes: bytes, mime: str = "image/png") -> str:
    """调用云端视觉模型识别图片中的文字内容（硅基流动 Qwen-VL）"""
    global last_image_error
    last_image_error = ""
    try:
        import config
        import requests as _requests
        b64 = base64.b64encode(file_bytes).decode("utf-8")
        url = f"{config.EMB_BASE_URL}/chat/completions"
        headers = {
            "Authorization": f"Bearer {config.EMB_API_KEY}",
            "Content-Type": "application/json"
        }
        image_block = {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}}
        text_block = {"type": "text", "text": "请完整识别并输出图片中的文字内容；如果图片没有文字，请用一两句话概括图片内容。"}
        # 按顺序尝试的视觉模型（需与硅基流动账号实际可用的模型一致）
        models = ["Qwen/Qwen3-VL-8B-Instruct"]
        last_error = ""
        for model in models:
            payload = {
                "model": model,
                "messages": [{"role": "user", "content": [image_block, text_block]}],
                "max_tokens": 2000
            }
            resp = _requests.post(url, headers=headers, json=payload, timeout=60)
            if resp.status_code != 200:
                last_error = _friendly_api_error(resp.status_code, resp.text)
                print(f"⚠️图片识别模型 {model} 失败：{last_error}")
                continue
            data = resp.json()
            content = data["choices"][0]["message"]["content"]
            if not content or not content.strip():
                last_error = "模型调用成功但未返回文字内容，请换一张更清晰的图片"
                print(f"⚠️图片识别模型 {model} 返回空内容")
                continue
            return clean_document_text(content)
        last_image_error = last_error
        print(f"⚠️图片识别全部模型失败，最后错误：{last_error}")
        # API 不可用时（如余额不足）自动尝试免费本地 OCR
        local_text = _ocr_image_bytes_local(file_bytes, mime)
        if local_text.strip():
            last_image_error = ""
            return local_text
        if last_local_ocr_error:
            last_image_error = last_local_ocr_error
        return ""
    except Exception as e:
        last_image_error = str(e)
        print(f"⚠️图片识别失败：{e}")
        return ""


def extract_file_text(filename: str, file_bytes: bytes):
    """按扩展名提取文本，返回 (是否支持, 文本)"""
    suffix = filename.lower()
    if suffix.endswith(".pdf"):
        return True, extract_pdf_bytes_text(file_bytes)
    if suffix.endswith(".docx"):
        return True, extract_docx_bytes_text(file_bytes)
    if suffix.endswith(".pptx"):
        return True, extract_pptx_bytes_text(file_bytes)
    if suffix.endswith(".doc"):
        return True, extract_old_doc_bytes_text(file_bytes)
    if suffix.endswith((".jpg", ".jpeg", ".png", ".bmp", ".webp")):
        mime_map = {
            ".jpg": "image/jpeg",
            ".jpeg": "image/jpeg",
            ".png": "image/png",
            ".bmp": "image/bmp",
            ".webp": "image/webp",
        }
        mime = mime_map.get(suffix, "image/png")
        return True, extract_image_text(file_bytes, mime)
    return False, ""
